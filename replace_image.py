from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def replace_screenshot():
    prs = Presentation("TDM_Nexus_Combined.pptx")
    screenshot_path = "actual_tdm_screenshot.png"
    replaced = False

    for slide in prs.slides:
        slide_has_target = False
        
        # Check all text shapes for the title text
        for shape in slide.shapes:
            if shape.has_text_frame:
                if "Feature Spotlight" in shape.text or "Centralized Delivery Hub" in shape.text:
                    slide_has_target = True
                    break
                    
        if slide_has_target:
            print("Found target slide.")
            
            # Find picture shape
            pic_shape = None
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    pic_shape = shape
                    break
            
            if pic_shape:
                left = pic_shape.left
                top = pic_shape.top
                width = pic_shape.width
                height = pic_shape.height
                
                # remove old shape
                sp = pic_shape._element
                sp.getparent().remove(sp)
                
                # add new picture
                slide.shapes.add_picture(screenshot_path, left, top, width, height)
                print("Replaced picture.")
                replaced = True
            else:
                print("No picture found on the slide.")

    if replaced:
        prs.save("TDM_Nexus_Combined_Final.pptx")
        print("Saved updated presentation as TDM_Nexus_Combined_Final.pptx")
    else:
        print("Could not find the slide or picture to replace.")

if __name__ == "__main__":
    replace_screenshot()
