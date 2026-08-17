import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create presentation
prs = Presentation()

def add_bullet(tf, text, level=0):
    p = tf.add_paragraph()
    p.text = text
    p.level = level
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(64, 64, 64)
    return p

# --- Title Slide ---
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "TDM Nexus\nCapabilities, Gains, & Reflections"
subtitle.text = "Professional Web Application for Technical Delivery Managers"
for p in title.text_frame.paragraphs:
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 153)

# --- Slide 1: Introduction to TDM Nexus ---
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Introduction to TDM Nexus"
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 153)

tf = content.text_frame
tf.text = "What is TDM Nexus?"
tf.paragraphs[0].font.size = Pt(28)
tf.paragraphs[0].font.bold = True

add_bullet(tf, "A comprehensive dashboard and management tool tailored for Technical Delivery Managers (TDMs).")
add_bullet(tf, "Centralized platform for tracking releases, finances, quality, and governance.")
add_bullet(tf, "Replaces fragmented tools with an interactive, responsive, and dynamic user interface.")
add_bullet(tf, "Features 3D visualizations and real-time tracking.")

# --- Slide 2: Core Capabilities (1/2) ---
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Core Capabilities (1/2)"
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 153)

tf = content.text_frame
tf.text = "Key features that drive efficiency:"
tf.paragraphs[0].font.size = Pt(28)
tf.paragraphs[0].font.bold = True

add_bullet(tf, "Release Planning & Governance:")
add_bullet(tf, "Dedicated modules for Release Planning Meetings (RPM).", level=1)
add_bullet(tf, "POAP and full lifecycle tracking.", level=1)

add_bullet(tf, "Financial Management & Approvals:")
add_bullet(tf, "End-to-end tracking of budgets.", level=1)
add_bullet(tf, "Monitoring ICAR movements and domain cross-charges.", level=1)

# --- Slide 3: Core Capabilities (2/2) ---
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Core Capabilities (2/2)"
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 153)

tf = content.text_frame
tf.text = "Ensuring quality and seamless operation:"
tf.paragraphs[0].font.size = Pt(28)
tf.paragraphs[0].font.bold = True

add_bullet(tf, "Quality Assurance:")
add_bullet(tf, "Integrated defect tracking.", level=1)
add_bullet(tf, "Testing quality metrics.", level=1)
add_bullet(tf, "Hypercare (ELS) monitoring.", level=1)

add_bullet(tf, "Seamless Workflows:")
add_bullet(tf, "Contextual tools including walkthrough wizards.", level=1)
add_bullet(tf, "Implementation tracking and data analyzing dashboards.", level=1)


# --- Slide 4: User Experience & Design ---
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "User Experience & Design"
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 153)

tf = content.text_frame
tf.text = "Designed for modern workflows:"
tf.paragraphs[0].font.size = Pt(28)
tf.paragraphs[0].font.bold = True

add_bullet(tf, "Fully responsive design accommodating desktop, tablet, and mobile views to keep TDMs connected.")
add_bullet(tf, "Contextual views with adaptable layouts, split-modes, and data-rich tables.")
add_bullet(tf, "Immersive data visualization powered by ThreeJS 3D canvas.")
add_bullet(tf, "Optimized GPU performance dynamically based on the device.")
add_bullet(tf, "Theming support, intelligent breakpoints, and accessible touch targets.")


# --- Slide 5: Key Gains ---
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Key Gains from Using TDM Nexus"
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 153)

tf = content.text_frame
tf.text = "Measurable improvements in daily operations:"
tf.paragraphs[0].font.size = Pt(28)
tf.paragraphs[0].font.bold = True

add_bullet(tf, "Process Efficiency: Streamlined planning to delivery processes via a 'single pane of glass'.")
add_bullet(tf, "Real-Time Visibility: Immediate insight into RAG statuses (Red, Amber, Green).")
add_bullet(tf, "Platform Synchronization: Seamlessly bridges the gap between Planning and Release Management.")
add_bullet(tf, "Robust Cost Control: Enhanced visibility into VROMs, budgets, and domain timesheet tracking.")

# --- Slide 6: Reflections & Future Outlook ---
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Reflections & Future Outlook"
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 153)

tf = content.text_frame
tf.text = "Looking ahead:"
tf.paragraphs[0].font.size = Pt(28)
tf.paragraphs[0].font.bold = True

add_bullet(tf, "Continuous Improvement: The ongoing transition towards a fully mobile-friendly UX.")
add_bullet(tf, "Improved Data Quality: Enforcing data hygiene in interconnected ADO fields.")
add_bullet(tf, "Empowered Collaboration: Cross-domain integration enables a true DevOps and Helix model.")
add_bullet(tf, "Scalable Architecture: Easily adapts to new portfolios, expanded teams, and future capabilities.")

# Save presentation
prs.save("TDM_Nexus_Presentation_v2.pptx")
print("Presentation generated successfully: TDM_Nexus_Presentation_v2.pptx")
